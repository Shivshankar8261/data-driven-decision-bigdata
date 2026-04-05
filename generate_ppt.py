"""
PPT Generator v2 - Enhanced with embedded images
---------------------------------------------------
Fills 6 slides with content + embedded chart images and architecture diagram.
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = "BDCCT_PPT_Template (1).pptx"
OUTPUT = "BDCCT_PPT.pptx"
CHARTS = "charts"

prs = Presentation(TEMPLATE)

SLIDE_W = prs.slide_width   # 9144000 EMU = 10 inches
SLIDE_H = prs.slide_height  # 6858000 EMU = 7.5 inches


def clear_placeholder(placeholder):
    """Remove placeholder so we can place free-form content."""
    sp = placeholder._element
    sp.getparent().remove(sp)


def add_textbox(slide, left, top, width, height, lines, font_size=12,
                bold_map=None, color_map=None, alignment=PP_ALIGN.LEFT):
    """Add a textbox with multiple lines, optional per-line bold/color."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = alignment

        is_bullet = line.strip().startswith('\u2022')
        if is_bullet:
            para.level = 1

        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = 'Calibri'

        if bold_map and i in bold_map:
            run.font.bold = True
        if color_map and i in color_map:
            run.font.color.rgb = RGBColor(*color_map[i])

    return txBox


def add_image(slide, image_path, left, top, width=None, height=None):
    """Add an image to a slide at specified position."""
    if not os.path.exists(image_path):
        return None
    kwargs = {}
    if width:
        kwargs['width'] = Inches(width)
    if height:
        kwargs['height'] = Inches(height)
    return slide.shapes.add_picture(image_path, Inches(left), Inches(top), **kwargs)


# ============================================================================
# SLIDE 1: Project Title & Problem Context
# ============================================================================
slide1 = prs.slides[0]
title1 = slide1.placeholders[0]
content1 = slide1.placeholders[1]

title1.text_frame.paragraphs[0].runs[0].text = "Project Title & Problem Context" if title1.text_frame.paragraphs[0].runs else ""
title1.text_frame.paragraphs[0].text = "Project Title & Problem Context"

content1_tf = content1.text_frame
content1_tf.clear()

lines_s1 = [
    ("Project Title:", True, RGBColor(44, 62, 80), 16),
    ("Data-Driven Decision Making in an Organization", False, RGBColor(41, 128, 185), 14),
    ("Using Big Data Technologies", False, RGBColor(41, 128, 185), 14),
    ("", False, None, 10),
    ("Name & USN: ____________________ | ____________________", False, RGBColor(80, 80, 80), 12),
    ("Domain: Big Data Analytics & Organizational Decision Intelligence", False, RGBColor(80, 80, 80), 12),
    ("", False, None, 8),
    ("Problem Statement:", True, RGBColor(192, 57, 43), 13),
    ("Organizations generate vast operational data but lack integrated", False, None, 12),
    ("pipelines to convert it into high-quality decisions. Poor data quality,", False, None, 12),
    ("high error rates, and fragmented analytics reduce decision", False, None, 12),
    ("effectiveness by up to 40%.", False, None, 12),
    ("", False, None, 8),
    ("Target Users: C-suite executives | Department heads | Data analysts | HR managers", True, RGBColor(39, 174, 96), 11),
]

for i, (text, bold, color, size) in enumerate(lines_s1):
    para = content1_tf.paragraphs[0] if i == 0 else content1_tf.add_paragraph()
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = 'Calibri'
    if color:
        run.font.color.rgb = color


# ============================================================================
# SLIDE 2: System Architecture (with diagram image)
# ============================================================================
slide2 = prs.slides[1]
content2 = slide2.placeholders[1]
clear_placeholder(content2)

add_image(slide2, os.path.join(CHARTS, '00_system_architecture.png'),
          left=0.4, top=1.5, width=9.2)

add_textbox(slide2, 0.5, 6.5, 9.0, 0.8, [
    "Tools: PySpark | GCP Dataproc | Cloud Storage | Parquet | Streamlit | Plotly | Python"
], font_size=10, bold_map={0: True}, color_map={0: (100, 100, 100)})


# ============================================================================
# SLIDE 3: PySpark Implementation (split: text left, charts right)
# ============================================================================
slide3 = prs.slides[2]
content3 = slide3.placeholders[1]
clear_placeholder(content3)

add_textbox(slide3, 0.3, 1.5, 4.8, 5.5, [
    "Dataset: 123,847 rows \u00d7 15 columns",
    "Period: Jan 2024 \u2013 Dec 2025",
    "",
    "Key Operations:",
    "\u2022 Null Imputation (median)",
    "\u2022 Duplicate Removal",
    "\u2022 Outlier Capping (IQR method)",
    "\u2022 Feature Engineering (9 new features)",
    "\u2022 Dept/Region Aggregations",
    "",
    "Code Snippet:",
    "df = spark.read.csv(path, schema=schema)",
    "df = df.dropDuplicates()",
    "df = df.withColumn('Efficiency_Score',",
    "    col('Perf')/(col('Time')+1))",
    "df.write.partitionBy('Department')",
    "    .parquet('processed_data/')",
    "",
    "Output: 24 cols, Parquet, partitioned",
], font_size=10,
   bold_map={0: True, 1: True, 3: True, 10: True, 17: True},
   color_map={10: (41, 128, 185), 11: (41, 128, 185), 12: (41, 128, 185),
              13: (41, 128, 185), 14: (41, 128, 185), 15: (41, 128, 185), 16: (41, 128, 185)})

add_image(slide3, os.path.join(CHARTS, '01_correlation_heatmap.png'),
          left=5.3, top=1.5, width=4.4)

add_textbox(slide3, 5.5, 6.2, 4.0, 0.4, [
    "Correlation Matrix of Key Variables"
], font_size=8, color_map={0: (120, 120, 120)}, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 4: Dashboard & Insights (with charts)
# ============================================================================
slide4 = prs.slides[3]
content4 = slide4.placeholders[1]
clear_placeholder(content4)

add_image(slide4, os.path.join(CHARTS, '02_dept_decision_impact.png'),
          left=0.3, top=1.4, width=4.5)

add_image(slide4, os.path.join(CHARTS, '04_monthly_trend.png'),
          left=5.0, top=1.4, width=4.7)

add_textbox(slide4, 0.3, 4.6, 9.3, 2.8, [
    "Key Business Insights:",
    "",
    "\u2022 Insight 1: Data Quality is the strongest predictor of Decision Impact (r = 0.736)",
    "\u2022 Insight 2: Training >40 hrs/yr yields 27% higher performance (4.76 vs 3.76)",
    "\u2022 Insight 3: Error rates >15% cause 43% drop in Decision Impact (77.7 \u2192 44.6)",
    "\u2022 Insight 4: Engineering leads Impact (+40% vs HR) \u2014 best data quality + training",
    "\u2022 Insight 5: Attrition Risk \u2194 Satisfaction feedback loop (r = \u20130.875)",
], font_size=11,
   bold_map={0: True},
   color_map={0: (44, 62, 80)})


# ============================================================================
# SLIDE 5: Real-World Application (with scatter + pie)
# ============================================================================
slide5 = prs.slides[4]
content5 = slide5.placeholders[1]
clear_placeholder(content5)

add_textbox(slide5, 0.3, 1.4, 4.6, 3.0, [
    "Real-World Applications:",
    "",
    "\u2022 Consulting \u2014 Client decision quality analysis",
    "\u2022 Healthcare \u2014 Clinical error monitoring",
    "\u2022 Retail \u2014 Regional decision optimization",
    "\u2022 Finance \u2014 Error-rate quality gates",
    "",
    "Scalability:",
    "\u2022 GCP Dataproc: 10\u2013100 node clusters",
    "\u2022 Parquet + Snappy: 60\u201380% compression",
    "\u2022 Delta Lake for ACID transactions",
    "",
    "Future:",
    "\u2022 Spark Structured Streaming (real-time)",
    "\u2022 ML models: predict Decision Impact",
    "\u2022 Automated threshold alerts",
], font_size=10,
   bold_map={0: True, 7: True, 12: True},
   color_map={0: (44, 62, 80), 7: (44, 62, 80), 12: (44, 62, 80)})

add_image(slide5, os.path.join(CHARTS, '03_quality_vs_impact_scatter.png'),
          left=5.1, top=1.4, width=4.6, height=2.8)

add_image(slide5, os.path.join(CHARTS, '06_decision_type_pie.png'),
          left=5.6, top=4.4, width=3.5, height=2.8)


# ============================================================================
# SLIDE 6: Conclusion & Questions
# ============================================================================
slide6 = prs.slides[5]
content6 = slide6.placeholders[1]
content6_tf = content6.text_frame
content6_tf.clear()

lines_s6 = [
    ("Summary:", True, RGBColor(44, 62, 80), 16),
    ("", False, None, 8),
    ("Problem:", True, RGBColor(192, 57, 43), 13),
    ("Organizations fail to leverage data for high-quality decisions", False, None, 12),
    ("due to poor data quality, fragmented pipelines, and lack of insights.", False, None, 12),
    ("", False, None, 8),
    ("Solution:", True, RGBColor(41, 128, 185), 13),
    ("Complete big data pipeline: PySpark + GCP + Streamlit", False, None, 12),
    ("123,847 records \u2192 Clean \u2192 Transform \u2192 Parquet \u2192 Dashboard", False, None, 12),
    ("", False, None, 8),
    ("Outcome:", True, RGBColor(39, 174, 96), 13),
    ("Data Quality (r=0.74) and Error Rate (r=\u20130.77) identified as top drivers.", False, None, 12),
    ("5 actionable decisions projected to improve quality by 18\u201325%.", False, None, 12),
    ("", False, None, 10),
    ("", False, None, 10),
    ("Thank You", True, RGBColor(44, 62, 80), 20),
]

for i, (text, bold, color, size) in enumerate(lines_s6):
    para = content6_tf.paragraphs[0] if i == 0 else content6_tf.add_paragraph()
    if text == "Thank You":
        para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = 'Calibri'
    if color:
        run.font.color.rgb = color


# ============================================================================
# SAVE
# ============================================================================
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
print("Embedded images: architecture diagram (slide 2), correlation heatmap (slide 3),")
print("  dept impact + monthly trend (slide 4), scatter + pie (slide 5)")
